"""
Basic tests for ESG Factsheet AI API.
Run with: pytest test_main.py
"""
import io
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_utils import (
    extract_text_from_pptx,
    find_ai_summary_shape,
    insert_text_into_ai_summary,
    truncate_text_for_llm
)
from ai_service import format_summary_for_pptx, format_bullets_as_text


def create_test_pptx_bytes() -> bytes:
    """Create a minimal test PPTX in memory."""
    prs = Presentation()
    
    # Add a slide with some text
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Test ESG Factsheet"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Carbon emissions reduced by 20%"
    
    # Add AI_SUMMARY placeholder
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    summary_box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    summary_box.name = "AI_SUMMARY"
    summary_box.text_frame.text = "Placeholder text"
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def test_extract_text_from_pptx():
    """Test PPTX text extraction."""
    pptx_bytes = create_test_pptx_bytes()
    text = extract_text_from_pptx(pptx_bytes)
    
    assert "Test ESG Factsheet" in text
    assert "Carbon emissions" in text
    assert "Slide 1" in text


def test_find_ai_summary_shape():
    """Test finding AI_SUMMARY shape."""
    pptx_bytes = create_test_pptx_bytes()
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    result = find_ai_summary_shape(prs)
    assert result is not None
    
    slide_idx, shape = result
    assert slide_idx == 1  # Second slide (0-indexed)
    assert shape.name == "AI_SUMMARY"


def test_insert_text_into_ai_summary():
    """Test inserting text into AI_SUMMARY placeholder."""
    pptx_bytes = create_test_pptx_bytes()
    
    summary_text = """**Strengths**
- Good carbon reduction performance
- Strong leadership commitment

**Weaknesses**
- Limited water management disclosure

**Action Plan (12 months)**
- Develop comprehensive water strategy
- Enhance ESG reporting framework"""
    
    regenerated_bytes = insert_text_into_ai_summary(pptx_bytes, summary_text)
    
    # Verify the text was inserted
    prs = Presentation(io.BytesIO(regenerated_bytes))
    result = find_ai_summary_shape(prs)
    assert result is not None
    
    _, shape = result
    assert "Strengths" in shape.text_frame.text
    assert "carbon reduction" in shape.text_frame.text


def test_insert_text_without_ai_summary():
    """Test error when AI_SUMMARY shape is missing."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "No AI_SUMMARY here"
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    pptx_bytes = output.read()
    
    with pytest.raises(ValueError, match="AI_SUMMARY shape not found"):
        insert_text_into_ai_summary(pptx_bytes, "Test text")


def test_truncate_text_for_llm():
    """Test text truncation for LLM context limits."""
    short_text = "This is a short text"
    assert truncate_text_for_llm(short_text, max_chars=100) == short_text
    
    long_text = "A" * 100000
    truncated = truncate_text_for_llm(long_text, max_chars=1000)
    assert len(truncated) <= 1100  # Allow some buffer for truncation notice
    assert "truncated" in truncated.lower()


def test_format_summary_for_pptx():
    """Test formatting summary bullets for PPTX."""
    strengths = ["Strong ESG governance", "Good emissions performance"]
    weaknesses = ["Limited water disclosure"]
    action_plan = ["Develop water strategy", "Enhance reporting"]
    
    formatted = format_summary_for_pptx(strengths, weaknesses, action_plan)
    
    assert "**Strengths**" in formatted
    assert "**Weaknesses**" in formatted
    assert "**Action Plan (12 months)**" in formatted
    assert "- Strong ESG governance" in formatted
    assert "- Limited water disclosure" in formatted


def test_format_bullets_as_text():
    """Test formatting list of bullets as text."""
    bullets = ["First bullet", "Second bullet", "Third bullet"]
    formatted = format_bullets_as_text(bullets)
    
    assert formatted == "- First bullet\n- Second bullet\n- Third bullet"


# Health check test (requires FastAPI TestClient)
def test_health_check():
    """Test the health check endpoint."""
    from fastapi.testclient import TestClient
    from main import app
    
    client = TestClient(app)
    response = client.get("/healthz")
    
    assert response.status_code == 200
    assert response.json()["ok"] is True


if __name__ == "__main__":
    pytest.main([__file__, "-v"])

