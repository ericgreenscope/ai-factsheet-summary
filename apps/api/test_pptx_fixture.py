"""
Generate a test PPTX file with an AI_SUMMARY placeholder.
Run this script to create a test fixture for development and testing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt


def create_test_pptx(filename: str = "test_factsheet.pptx"):
    """
    Create a test ESG factsheet PPTX with AI_SUMMARY placeholder.
    
    Args:
        filename: Output filename
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "ESG Factsheet 2024"
    subtitle.text = "Sample Company Inc. - Sustainability Assessment"
    
    # Slide 2: Content with ESG data
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "Environmental Performance"
    
    body_shape = slide2.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Key Metrics:"
    
    p = tf.add_paragraph()
    p.text = "Carbon emissions reduced by 15% year-over-year"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Renewable energy usage increased to 45%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Water consumption decreased by 12%"
    p.level = 1
    
    # Slide 3: Social metrics
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    title.text = "Social Impact"
    
    body_shape = slide3.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Workforce & Community:"
    
    p = tf.add_paragraph()
    p.text = "Employee satisfaction score: 78/100"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Diversity: 40% women in leadership roles"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Community investment: $2M in local programs"
    p.level = 1
    
    # Slide 4: Governance
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    title.text = "Governance"
    
    body_shape = slide4.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Corporate Governance:"
    
    p = tf.add_paragraph()
    p.text = "Independent board members: 60%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "ESG committee established in 2023"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Annual sustainability reporting since 2020"
    p.level = 1
    
    # Slide 5: AI Summary Placeholder (THIS IS THE KEY SLIDE)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide5.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "AI-Generated Executive Summary"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    
    # Add AI_SUMMARY text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    summary_box = slide5.shapes.add_textbox(left, top, width, height)
    summary_box.name = "AI_SUMMARY"  # This is the key identifier
    
    text_frame = summary_box.text_frame
    text_frame.text = "[This space will be filled with AI-generated summary]"
    text_frame.word_wrap = True
    
    # Save presentation
    prs.save(filename)
    print(f"Test PPTX created: {filename}")
    print(f"  - 5 slides created")
    print(f"  - AI_SUMMARY placeholder on Slide 5")
    print(f"  - Ready for testing!")


if __name__ == "__main__":
    create_test_pptx()

