"""PDF conversion utilities for PPTX files."""
import io
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image
import tempfile
from pathlib import Path


def convert_pptx_to_pdf(pptx_bytes: bytes) -> bytes:
    """
    Convert PPTX file to PDF using Python libraries.

    This creates a simple PDF representation of the PPTX slides.
    For better quality, consider using a cloud conversion service.

    Args:
        pptx_bytes: PPTX file content as bytes

    Returns:
        PDF file content as bytes

    Raises:
        RuntimeError: If conversion fails
    """
    try:
        # Load presentation
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)
        width, height = letter

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract text from slide
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)

            # Draw slide number
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, f"Slide {slide_num}")

            # Draw text content
            c.setFont("Helvetica", 12)
            y_position = height - 100
            for text in text_content:
                # Split long text into lines
                lines = text.split('\n')
                for line in lines:
                    if y_position < 50:  # Start new page if needed
                        c.showPage()
                        y_position = height - 50

                    # Truncate very long lines
                    if len(line) > 80:
                        line = line[:80] + "..."

                    c.drawString(50, y_position, line)
                    y_position -= 20

            # Add page for next slide
            if slide_num < len(prs.slides):
                c.showPage()

        # Save PDF
        c.save()
        pdf_buffer.seek(0)
        return pdf_buffer.read()

    except Exception as e:
        raise RuntimeError(f"PDF conversion failed: {str(e)}")


def is_conversion_available() -> bool:
    """
    Check if PDF conversion is available.

    Returns:
        True (always available with Python libraries)
    """
    return True
