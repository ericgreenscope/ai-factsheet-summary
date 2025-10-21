"""PPTX text extraction and AI_SUMMARY shape manipulation."""
import io
from typing import Optional, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from markdown_it import MarkdownIt
from markdown_it.token import Token


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """
    Extract all text from a PPTX file.
    
    Args:
        pptx_bytes: Binary PPTX data
    
    Returns:
        Flattened text from all slides
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    all_text = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
            
            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text.append(cell.text.strip())
        
        if slide_text:
            all_text.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text))
    
    return "\n\n".join(all_text)


def find_ai_summary_shape(prs: Presentation) -> Optional[Tuple[int, object]]:
    """
    Find the shape with name or alt_text = "AI_SUMMARY".
    
    Args:
        prs: PowerPoint presentation object
    
    Returns:
        Tuple of (slide_index, shape) or None if not found
    """
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check shape name
            if hasattr(shape, "name") and shape.name == "AI_SUMMARY":
                return (slide_idx, shape)
            
            # Check alt text
            if hasattr(shape, "element") and hasattr(shape.element, "get"):
                alt_text = shape.element.get("descr", "")
                if alt_text == "AI_SUMMARY":
                    return (slide_idx, shape)
    
    return None


def insert_text_into_ai_summary(
    pptx_bytes: bytes,
    summary_text: str
) -> bytes:
    """
    Insert summary text into the AI_SUMMARY placeholder shape.
    
    Args:
        pptx_bytes: Original PPTX binary data
        summary_text: Formatted summary text to insert
    
    Returns:
        Modified PPTX binary data
    
    Raises:
        ValueError: If AI_SUMMARY shape not found
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # Find the AI_SUMMARY shape
    result = find_ai_summary_shape(prs)
    if not result:
        raise ValueError("AI_SUMMARY shape not found in presentation")
    
    slide_idx, shape = result
    
    # Check if shape has text frame
    if not hasattr(shape, "text_frame"):
        raise ValueError("AI_SUMMARY shape does not support text")
    
    # Clear existing text and insert new summary
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Add the summary text
    p = text_frame.paragraphs[0]
    p.text = summary_text
    p.font.size = Pt(11)  # Default font size
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.read()


def truncate_text_for_llm(text: str, max_chars: int = 80000) -> str:
    """
    Truncate extracted text to fit within LLM context limits.
    
    Args:
        text: Full extracted text
        max_chars: Maximum character count (default ~80k for safety)
    
    Returns:
        Truncated text
    """
    if len(text) <= max_chars:
        return text
    
    # Truncate and add notice
    truncated = text[:max_chars]
    truncated += "\n\n[... Text truncated to fit context limit ...]"
    
    return truncated


def insert_markdown_into_ai_summary(
    pptx_bytes: bytes,
    markdown_text: str
) -> bytes:
    """
    Parse Markdown and insert formatted text into AI_SUMMARY placeholder.
    
    Args:
        pptx_bytes: Original PPTX binary data
        markdown_text: Markdown formatted analysis text
    
    Returns:
        Modified PPTX binary data
    
    Raises:
        ValueError: If AI_SUMMARY shape not found
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # Find the AI_SUMMARY shape
    result = find_ai_summary_shape(prs)
    if not result:
        raise ValueError("AI_SUMMARY shape not found in presentation")
    
    slide_idx, shape = result
    
    # Check if shape has text frame
    if not hasattr(shape, "text_frame"):
        raise ValueError("AI_SUMMARY shape does not support text")
    
    # Parse Markdown
    md = MarkdownIt()
    tokens = md.parse(markdown_text)
    
    # Clear existing text
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Process tokens
    in_list = False
    is_ordered_list = False
    list_counter = 0
    
    i = 0
    while i < len(tokens):
        token = tokens[i]
        
        # Headings: h1/h2/h3 -> larger bold text
        if token.type == "heading_open":
            level = int(token.tag[1])
            if i + 1 < len(tokens) and tokens[i + 1].type == "inline":
                inline_token = tokens[i + 1]
                p = text_frame.add_paragraph()
                p.level = 0
                
                if level == 1:
                    p.font.size = Pt(24)
                elif level == 2:
                    p.font.size = Pt(20)
                else:
                    p.font.size = Pt(16)
                p.font.bold = True
                
                # Apply inline formatting to heading
                _apply_inline_formatting(p, inline_token)
                
                i += 3  # Skip heading_close
                continue
        
        # Bullet lists
        if token.type == "bullet_list_open":
            in_list = True
            is_ordered_list = False
            list_counter = 0
            i += 1
            continue
        
        if token.type == "ordered_list_open":
            in_list = True
            is_ordered_list = True
            list_counter = 0
            i += 1
            continue
        
        if token.type in ("bullet_list_close", "ordered_list_close"):
            in_list = False
            i += 1
            continue
        
        # List items
        if token.type == "list_item_open":
            i += 1
            continue
        
        if token.type == "list_item_close":
            i += 1
            continue
        
        # Paragraphs
        if token.type == "paragraph_open":
            if i + 1 < len(tokens) and tokens[i + 1].type == "inline":
                inline_token = tokens[i + 1]
                p = text_frame.add_paragraph()
                p.level = 0
                p.font.size = Pt(12)
                
                # Add list prefix if needed
                list_prefix = ""
                if in_list:
                    list_counter += 1
                    if is_ordered_list:
                        list_prefix = f"{list_counter}. "
                    else:
                        list_prefix = "• "
                
                # Add prefix as plain text first, then apply inline formatting
                if list_prefix:
                    run = p.add_run()
                    run.text = list_prefix
                
                # Process inline formatting
                _apply_inline_formatting(p, inline_token)
                
                i += 2
                continue
        
        if token.type == "paragraph_close":
            i += 1
            continue
        
        i += 1
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output.read()


def _apply_inline_formatting(paragraph, inline_token):
    """
    Apply formatting (bold, italic, code) to inline content in a paragraph.
    Properly handles markdown-it token children.
    
    Args:
        paragraph: python-pptx paragraph object
        inline_token: Markdown-it inline token with children
    """
    if not inline_token.children:
        # No children, just add the raw content
        run = paragraph.add_run()
        run.text = inline_token.content
        return
    
    # Track formatting state
    is_bold = False
    is_italic = False
    
    i = 0
    children = inline_token.children
    while i < len(children):
        child = children[i]
        
        if child.type == "text":
            run = paragraph.add_run()
            run.text = child.content
            if is_bold:
                run.font.bold = True
            if is_italic:
                run.font.italic = True
        
        elif child.type == "strong_open":
            is_bold = True
        
        elif child.type == "strong_close":
            is_bold = False
        
        elif child.type == "em_open":
            is_italic = True
        
        elif child.type == "em_close":
            is_italic = False
        
        elif child.type == "code_inline":
            run = paragraph.add_run()
            run.text = child.content
            run.font.name = "Consolas"
            run.font.size = Pt(10)
            if is_bold:
                run.font.bold = True
            if is_italic:
                run.font.italic = True
        
        elif child.type in ("link_open", "link_close", "softbreak", "hardbreak"):
            # Links: render text only (no URL)
            # Breaks: skip
            pass
        
        i += 1

